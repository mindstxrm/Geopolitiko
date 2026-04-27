"""Generate Scenario Engine exports as DOCX and PPTX; also custom briefing DOCX."""
from io import BytesIO


def build_briefing_docx(
    title: str,
    intro: str,
    articles: list,
    format_type: str = "summary_link",
    grouped: list = None,
    executive_summary: str = None,
    report_date: str = None,
    sensitivity_tier: str = "internal",
    legal_review_required: bool = False,
) -> bytes:
    """Build a DOCX from briefing title, intro, and articles. Supports all PDF format types and grouping."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    valid_formats = ("minimal", "summary_only", "summary_link", "with_takeaways", "executive", "full")
    if format_type not in valid_formats:
        format_type = "summary_link"
    tier = (sensitivity_tier or "internal").strip().lower()
    if tier not in ("public", "internal", "restricted"):
        tier = "internal"
    tier_label = {"public": "PUBLIC", "internal": "INTERNAL USE ONLY", "restricted": "RESTRICTED"}[tier]

    doc = Document()
    section = doc.sections[0]
    header = section.header
    if header.paragraphs:
        hp = header.paragraphs[0]
        hp.text = ""
    else:
        hp = header.add_paragraph()
    r0 = hp.add_run(f"{tier_label} · Geopolitical Terminal")
    r0.font.size = Pt(9)
    r0.font.color.rgb = RGBColor(120, 120, 120)
    if report_date:
        r1 = hp.add_run(f" · {report_date}")
        r1.font.size = Pt(9)
        r1.font.color.rgb = RGBColor(120, 120, 120)
    if legal_review_required:
        hp2 = header.add_paragraph()
        r2 = hp2.add_run("Legal review required before external distribution.")
        r2.font.size = Pt(8)
        r2.font.color.rgb = RGBColor(180, 80, 80)
        r2.bold = True

    doc.add_heading(title or "Briefing", 0)
    if report_date:
        p = doc.add_paragraph()
        p.add_run(f"Report date: {report_date}").italic = True
    if intro:
        doc.add_paragraph(intro)
    if executive_summary:
        doc.add_heading("Executive summary", level=1)
        doc.add_paragraph(executive_summary)
    doc.add_paragraph()

    groups = grouped if grouped else [("", articles or [])]
    idx = 1
    for group_label, group_articles in groups:
        if group_label:
            doc.add_heading(group_label, level=1)
        for a in group_articles:
            imp = a.get("impact_score")
            imp_str = f"Impact {imp}/10" if imp is not None else ""
            date_str = (a.get("published_utc") or a.get("scraped_at") or "")[:10]
            meta_parts = [a.get("source_name") or "—"]
            if date_str:
                meta_parts.append(date_str)
            if imp_str:
                meta_parts.append(imp_str)
            if a.get("event_type"):
                meta_parts.append(a.get("event_type"))

            doc.add_paragraph(f"{idx}. {a.get('title') or 'Untitled'}", style="List Number")
            doc.add_paragraph(" · ".join(meta_parts), style="Intense Quote")
            if a.get("topics") and format_type in ("full", "executive"):
                doc.add_paragraph(f"Topics: {(a.get('topics') or '')[:100]}")

            if format_type == "minimal":
                if a.get("url"):
                    doc.add_paragraph(f"Link: {a['url']}")
            elif format_type == "summary_only":
                doc.add_paragraph((a.get("summary") or "")[:500] + ("…" if len(a.get("summary") or "") > 500 else ""))
            elif format_type == "executive":
                doc.add_paragraph((a.get("summary") or "")[:300] + ("…" if len(a.get("summary") or "") > 300 else ""))
                if a.get("url"):
                    doc.add_paragraph(f"Link: {a['url']}")
            elif format_type == "with_takeaways":
                doc.add_paragraph((a.get("summary") or "")[:400] + ("…" if len(a.get("summary") or "") > 400 else ""))
                if a.get("key_takeaways"):
                    doc.add_paragraph(f"Key takeaways: {(a.get('key_takeaways') or '')[:400]}")
                if a.get("url"):
                    doc.add_paragraph(f"Link: {a['url']}")
            elif format_type == "full":
                doc.add_paragraph((a.get("summary") or "")[:400] + ("…" if len(a.get("summary") or "") > 400 else ""))
                if a.get("why_it_matters"):
                    doc.add_paragraph(f"Why it matters: {(a.get('why_it_matters') or '')[:300]}")
                if a.get("key_takeaways"):
                    doc.add_paragraph(f"Key takeaways: {(a.get('key_takeaways') or '')[:300]}")
                if a.get("url"):
                    doc.add_paragraph(f"Link: {a['url']}")
            else:
                doc.add_paragraph((a.get("summary") or "")[:400] + ("…" if len(a.get("summary") or "") > 400 else ""))
                if a.get("url"):
                    doc.add_paragraph(f"Link: {a['url']}")
            doc.add_paragraph()
            idx += 1

    doc.add_paragraph()
    footer_p = doc.add_paragraph()
    footer_p.add_run("Generated by Geopolitical Terminal").italic = True
    footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fd = doc.add_paragraph()
    fd.add_run(
        "Research support only — not legal or investment advice. Verify against primary sources."
    ).italic = True
    fd.alignment = WD_ALIGN_PARAGRAPH.CENTER

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def _agent_title(key: str) -> str:
    if key == "domestic_politics" or key == "political":
        return "Political agent"
    if key == "economic_policy":
        return "Economic agent"
    if key == "private_sector":
        return "Private sector response agent"
    if key == "military":
        return "Military agent"
    if key == "diplomatic":
        return "Diplomatic agent"
    return key.replace("_", " ").title()


def build_policy_memo_docx(engine_result: dict) -> bytes:
    from docx import Document

    event = engine_result.get("event_label") or engine_result.get("event_type") or "Scenario"
    region = engine_result.get("region") or ""
    country = engine_result.get("country") or ""
    run_at = (engine_result.get("run_at") or "")[:19].replace("T", " ")
    agents = engine_result.get("agents") or {}
    paths = engine_result.get("paths") or {}
    path_desc = engine_result.get("path_descriptions") or {}

    doc = Document()
    doc.add_heading(f"Policy memo: {event}", 0)
    doc.add_paragraph(f"Generated {run_at}")
    if region or country:
        doc.add_paragraph(f"Region: {region}" + (f" | Country: {country}" if country else ""))
    doc.add_paragraph()

    doc.add_heading("Summary", level=1)
    doc.add_paragraph(
        f"Multi-agent scenario simulation for {event}. Probability-weighted pathways and agent "
        "assessments below. For official use; update assumptions and re-run as conditions change."
    )
    doc.add_heading("Agent assessments", level=1)
    for key, text in agents.items():
        doc.add_heading(_agent_title(key), level=2)
        doc.add_paragraph(text)
    doc.add_heading("Outcome pathways", level=1)
    for key in ("contained", "regional_escalation", "systemic_crisis"):
        label = {"contained": "Contained", "regional_escalation": "Regional escalation", "systemic_crisis": "Systemic crisis"}.get(key, key)
        pct = paths.get(key, 0)
        desc = path_desc.get(key) or ""
        doc.add_paragraph(f"{label} ({pct}%): {desc}", style="List Bullet")
    doc.add_paragraph()
    doc.add_paragraph("Geopolitical Terminal · Scenario Engine · Confidential.")

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def build_risk_briefing_docx(engine_result: dict) -> bytes:
    from docx import Document

    event = engine_result.get("event_label") or engine_result.get("event_type") or "Scenario"
    region = engine_result.get("region") or ""
    country = engine_result.get("country") or ""
    run_at = (engine_result.get("run_at") or "")[:19].replace("T", " ")
    agents = engine_result.get("agents") or {}
    paths = engine_result.get("paths") or {}
    path_desc = engine_result.get("path_descriptions") or {}

    doc = Document()
    doc.add_heading(f"Risk briefing: {event}", 0)
    doc.add_paragraph(f"{run_at}")
    if region or country:
        doc.add_paragraph(f"Scope: {region}" + (f" — {country}" if country else ""))
    doc.add_paragraph()

    doc.add_heading("Risk scenario", level=1)
    doc.add_paragraph(f"Event type: {event}. The following pathways and agent views support risk appetite and contingency planning.")
    doc.add_heading("Probability-weighted pathways", level=1)
    for key in ("contained", "regional_escalation", "systemic_crisis"):
        label = {"contained": "Contained", "regional_escalation": "Regional escalation", "systemic_crisis": "Systemic crisis"}.get(key, key)
        pct = paths.get(key, 0)
        desc = path_desc.get(key) or ""
        doc.add_paragraph(f"{label} ({pct}%): {desc}", style="List Bullet")
    doc.add_heading("Multi-agent analysis", level=1)
    for key, text in agents.items():
        doc.add_heading(_agent_title(key), level=2)
        doc.add_paragraph(text)
    doc.add_heading("Recommendations", level=1)
    doc.add_paragraph("Monitor triggers and leading indicators for pathway shifts.", style="List Bullet")
    doc.add_paragraph("Align contingency plans with Contained and Regional escalation cases.", style="List Bullet")
    doc.add_paragraph("Stress-test exposures against Systemic crisis assumptions.", style="List Bullet")
    doc.add_paragraph()
    doc.add_paragraph("Geopolitical Terminal · Scenario Engine · Risk Briefing.")

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def build_investor_note_docx(engine_result: dict) -> bytes:
    from docx import Document

    event = engine_result.get("event_label") or engine_result.get("event_type") or "Scenario"
    region = engine_result.get("region") or ""
    country = engine_result.get("country") or ""
    run_at = (engine_result.get("run_at") or "")[:19].replace("T", " ")
    agents = engine_result.get("agents") or {}
    paths = engine_result.get("paths") or {}
    path_desc = engine_result.get("path_descriptions") or {}

    doc = Document()
    doc.add_heading(f"Investor note: {event}", 0)
    doc.add_paragraph(f"Scenario Engine output · {run_at}")
    if region or country:
        doc.add_paragraph(f"Region: {region}" + (f" | {country}" if country else ""))
    doc.add_paragraph()

    doc.add_heading("Event and pathways", level=1)
    doc.add_paragraph(f"{event} — Three-path probability distribution:")
    for key in ("contained", "regional_escalation", "systemic_crisis"):
        label = {"contained": "Contained", "regional_escalation": "Regional escalation", "systemic_crisis": "Systemic crisis"}.get(key, key)
        pct = paths.get(key, 0)
        desc = path_desc.get(key) or ""
        doc.add_paragraph(f"{label} ({pct}%): {desc}", style="List Bullet")
    doc.add_heading("Economic and private sector implications", level=1)
    for key, text in agents.items():
        doc.add_heading(_agent_title(key), level=2)
        doc.add_paragraph(text)
    doc.add_heading("Disclosure", level=1)
    doc.add_paragraph(
        "This note is scenario-based analysis, not a forecast. Probabilities are model-derived and subject to change. Not investment advice."
    )
    doc.add_paragraph()
    doc.add_paragraph("Geopolitical Terminal · Scenario Engine · Investor Note.")

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def build_exec_summary_pptx(engine_result: dict) -> bytes:
    from pptx import Presentation

    event = engine_result.get("event_label") or engine_result.get("event_type") or "Scenario"
    region = engine_result.get("region") or ""
    country = engine_result.get("country") or ""
    run_at = (engine_result.get("run_at") or "")[:19].replace("T", " ")
    agents = engine_result.get("agents") or {}
    paths = engine_result.get("paths") or {}
    path_desc = engine_result.get("path_descriptions") or {}

    prs = Presentation()
    # Title slide
    title_slide = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide)
    slide.shapes.title.text = f"Executive summary: {event}"
    slide.placeholders[1].text = f"Scenario Engine · {run_at}\nScope: {region or 'Global'}{f' — {country}' if country else ''}"

    # Three-path outcomes
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Three-path outcomes"
    body = slide.placeholders[1].text_frame
    for key in ("contained", "regional_escalation", "systemic_crisis"):
        label = {"contained": "Contained", "regional_escalation": "Regional escalation", "systemic_crisis": "Systemic crisis"}.get(key, key)
        pct = paths.get(key, 0)
        desc = (path_desc.get(key) or "")[:200]
        p = body.add_paragraph()
        p.text = f"{label} ({pct}%): {desc}"
        p.level = 0

    # Key agent takeaways
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Key agent takeaways"
    body = slide.placeholders[1].text_frame
    for key in ("economic_policy", "political", "military", "diplomatic", "private_sector"):
        title = _agent_title(key)
        text = (agents.get(key) or "")[:300]
        if text and len(text) >= 300:
            text = text[:297] + "..."
        p = body.add_paragraph()
        p.text = f"{title}: {text}"
        p.level = 0

    # Implications and next steps
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Implications and next steps"
    body = slide.placeholders[1].text_frame
    for line in [
        "Update monitoring and early-warning triggers.",
        "Align contingency plans with probability-weighted paths.",
        "Re-run scenario as new information arrives.",
    ]:
        p = body.add_paragraph()
        p.text = line
        p.level = 0

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
